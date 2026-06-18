"""
电商数据洞察平台 - 项目演示PPT生成脚本 v3
使用 python-pptx 生成专业演示文稿（22页）
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# 设计系统 - Corporate Professional 风格
# ============================================================
COLORS = {
    'primary': RGBColor(0x00, 0x33, 0x66),
    'secondary': RGBColor(0x00, 0x78, 0xD4),
    'accent': RGBColor(0x00, 0xB4, 0xD8),
    'dark': RGBColor(0x1B, 0x1B, 0x2F),
    'text': RGBColor(0x2D, 0x2D, 0x2D),
    'light_text': RGBColor(0xFF, 0xFF, 0xFF),
    'muted': RGBColor(0x6C, 0x75, 0x7D),
    'bg_light': RGBColor(0xF8, 0xF9, 0xFA),
    'green': RGBColor(0x28, 0xA7, 0x45),
    'orange': RGBColor(0xFD, 0x7E, 0x14),
    'red': RGBColor(0xDC, 0x35, 0x45),
}

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# 截图目录
PHOTOS_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'photos')


# ============================================================
# 工具函数
# ============================================================
def add_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=24,
                 color=None, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color or COLORS['text']
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=18,
                    color=None, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color or COLORS['text']
        p.font.name = 'Microsoft YaHei'
        p.space_after = spacing
    return txBox


def add_metric_card(slide, left, top, width, height, number, label, color=COLORS['secondary']):
    card = add_shape(slide, left, top, width, height, RGBColor(0xFF, 0xFF, 0xFF))
    add_shape(slide, left, top, width, Pt(4), color)
    add_text_box(slide, left, top + Pt(20), width, Inches(0.6),
                 str(number), font_size=32, color=color, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Pt(60), width, Inches(0.4),
                 label, font_size=14, color=COLORS['muted'],
                 alignment=PP_ALIGN.CENTER)


def add_section_divider(slide, title, subtitle=""):
    add_bg(slide, COLORS['primary'])
    add_shape(slide, Inches(1), Inches(3.2), Inches(2), Pt(3), COLORS['accent'])
    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
                 title, font_size=44, color=COLORS['light_text'], bold=True)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
                     subtitle, font_size=20, color=COLORS['accent'])


def add_flow_arrow(slide, left, top, width=Inches(0.8)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['secondary']
    shape.line.fill.background()


def add_header(slide, title):
    """添加统一的页面标题栏"""
    add_shape(slide, 0, 0, SLIDE_WIDTH, Inches(1.1), COLORS['primary'])
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
                 title, font_size=36, color=COLORS['light_text'], bold=True)


def add_screenshot(slide, filename, left, top, width, height=None):
    """嵌入截图，自动处理路径和边框"""
    filepath = os.path.join(PHOTOS_DIR, filename)
    if not os.path.exists(filepath):
        add_text_box(slide, left, top, width, height or Inches(3),
                     f"[截图缺失: {filename}]", font_size=14, color=COLORS['muted'])
        return
    # 添加白色背景框
    border = add_shape(slide, left - Pt(2), top - Pt(2),
                       width + Pt(4), (height or Inches(3)) + Pt(4),
                       RGBColor(0xDD, 0xDD, 0xDD))
    pic = slide.shapes.add_picture(filepath, left, top, width, height)
    # 如果没指定高度，按比例计算
    if height is None:
        # 已经添加了，python-pptx 会自动计算
        pass


def add_content_slide(slide, title):
    """创建内容页（浅色背景+标题栏）"""
    add_bg(slide, COLORS['bg_light'])
    add_header(slide, title)


# ============================================================
# 主生成逻辑
# ============================================================
def generate_ppt():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # ================================================================
    # Slide 1: 封面
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, COLORS['dark'])
    add_shape(slide, 0, 0, SLIDE_WIDTH, Pt(4), COLORS['accent'])
    add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
                 "电商数据洞察平台", font_size=54, color=COLORS['light_text'],
                 bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(3.1), Inches(10), Inches(0.8),
                 "基于 Spark + HBase + Airflow 的大数据分析系统",
                 font_size=24, color=COLORS['accent'], alignment=PP_ALIGN.CENTER)
    add_shape(slide, Inches(5), Inches(4.1), Inches(3), Pt(2), COLORS['accent'])
    # 团队信息（留空占位）
    add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.5),
                 "大数据技术课程项目", font_size=18, color=COLORS['muted'],
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.5),
                 "组员：___________    学号：___________    班级：___________",
                 font_size=16, color=COLORS['muted'], alignment=PP_ALIGN.CENTER)

    # ================================================================
    # Slide 2: 项目概述
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_content_slide(slide, "项目概述")

    cards = [
        ("2,000", "商品数量"),
        ("75.5万", "日销量记录"),
        ("22,970", "用户评价"),
        ("8", "商品品类"),
    ]
    for i, (num, label) in enumerate(cards):
        add_metric_card(slide, Inches(0.8 + i * 3.1), Inches(1.5),
                        Inches(2.7), Inches(1.2), num, label)

    add_text_box(slide, Inches(0.8), Inches(3.2), Inches(12), Inches(0.5),
                 "项目目标", font_size=24, color=COLORS['primary'], bold=True)
    goals = [
        "模拟抖音电商数据采集 → 自动化 ETL 管道 → 多维分析 → 可视化展示",
        "涵盖大数据全链路：Kafka 消息队列 → HDFS 分布式存储 → Spark 计算 → HBase 列式存储",
        "Apache Airflow 每日自动调度，实现数据管道全自动化运行",
        "支持爆款指数分析、NLP 情感分析、销量趋势预测等功能",
    ]
    add_bullet_text(slide, Inches(0.8), Inches(3.8), Inches(12), Inches(3),
                    goals, font_size=16, spacing=Pt(12))

    # ================================================================
    # Slide 3: 技术架构
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_content_slide(slide, "技术架构")

    layers = [
        ("调度层", "Apache Airflow", "DAG 工作流编排，每日凌晨自动执行", COLORS['secondary']),
        ("消息层", "Kafka", "数据传输通道，解耦数据生产与消费", COLORS['accent']),
        ("存储层", "HDFS + HBase + ES", "HDFS 分布式存储 | HBase 列式查询 | ES 全文搜索", COLORS['green']),
        ("计算层", "Spark (SQL + MLlib)", "数据清洗、商品分析、NLP 分析、趋势聚合", COLORS['orange']),
        ("服务层", "FastAPI + Vue3", "REST API 后端 | ECharts 可视化前端", COLORS['primary']),
    ]
    for i, (name, tech, desc, color) in enumerate(layers):
        y = Inches(1.4 + i * 1.1)
        add_shape(slide, Inches(0.8), y, Pt(6), Inches(0.8), color)
        add_text_box(slide, Inches(1.1), y, Inches(1.8), Inches(0.4),
                     name, font_size=16, color=COLORS['muted'], bold=True)
        add_text_box(slide, Inches(3), y, Inches(3.5), Inches(0.4),
                     tech, font_size=20, color=COLORS['primary'], bold=True)
        add_text_box(slide, Inches(6.5), y, Inches(6), Inches(0.8),
                     desc, font_size=14, color=COLORS['text'])

    # ================================================================
    # Slide 4: Airflow DAG 流程
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_content_slide(slide, "Airflow 自动化调度")

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(12), Inches(0.4),
                 "DAG: ecommerce_etl_pipeline  |  调度: 每天凌晨 2:00  |  全部在Docker容器内执行",
                 font_size=14, color=COLORS['muted'])

    steps = [
        ("1", "生成数据", "airflow容器内\n并行生成3种JSON", COLORS['secondary']),
        ("2", "初始化存储", "backend容器\n写入HBase+ES", COLORS['accent']),
        ("3", "上传HDFS", "hdfs容器\n原始数据存储", COLORS['green']),
        ("4", "Spark清洗", "spark容器\n去空/去重/格式化", COLORS['orange']),
        ("5", "Spark分析", "spark容器\n3个作业并行", COLORS['red']),
        ("6", "计算指标", "backend容器\n补充计算", COLORS['primary']),
    ]
    for i, (num, title, desc, color) in enumerate(steps):
        x = Inches(0.5 + i * 2.1)
        y = Inches(2.2)
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.6), y, Inches(0.6), Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.color.rgb = COLORS['light_text']
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        add_text_box(slide, x, y + Inches(0.8), Inches(1.8), Inches(0.4),
                     title, font_size=16, color=COLORS['primary'], bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, y + Inches(1.2), Inches(1.8), Inches(0.8),
                     desc, font_size=12, color=COLORS['muted'],
                     alignment=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_flow_arrow(slide, x + Inches(1.7), y + Inches(0.15), Inches(0.4))

    add_text_box(slide, Inches(0.8), Inches(4.8), Inches(12), Inches(0.4),
                 "任务依赖关系", font_size=18, color=COLORS['primary'], bold=True)
    deps = [
        "generate_* (并行) → init_hbase_es → create_hdfs_dirs → upload_* (并行) → spark_cleaning",
        "spark_cleaning → spark_product_analysis / spark_nlp_analysis / spark_trend_aggregation (并行)",
        "spark_analysis* → compute_analysis / compute_review_analysis → finish",
    ]
    add_bullet_text(slide, Inches(0.8), Inches(5.3), Inches(12), Inches(2),
                    deps, font_size=13, color=COLORS['text'], spacing=Pt(8))

    # ================================================================
    # Slide 5: Airflow 运行截图（新增）
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_content_slide(slide, "Airflow DAG 运行效果")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(12), Inches(0.4),
                 "DAG 全部任务执行成功（2026-06-17 运行记录）",
                 font_size=16, color=COLORS['muted'])
    add_screenshot(slide, 'airflow运行成功.png',
                   Inches(0.5), Inches(1.9), Inches(12.3), Inches(5.2))

    # ================================================================
    # Slide 6: 章节分隔 - 数据流转详解
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_section_divider(slide, "数据流转详解",
                        "从 Mock 数据生成到 Spark 计算，再到最终存储和前端展示")

    # ================================================================
    # Slide 7: Mock 数据生成
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_content_slide(slide, "阶段 1：Mock 数据生成 (Airflow 容器内)")

    data_sources = [
        ("商品数据", "products.json", "2,000 条", "8个品类 × 250商品\n品类: 手机/电脑/服装/美妆/食品/家居/数码/运动",
         "product_id, name, category, brand, price, original_price, description, shop_name"),
        ("评价数据", "reviews.json", "22,970 条", "每商品 8-15 条评价\n评分分布: 好评60% + 中评20% + 差评20%",
         "review_id, product_id, content, rating, sentiment, likes"),
        ("日销量数据", "monthly_sales.json", "730,000 条", "2000商品 × 365天\n含星期效应、月度趋势、大促活动",
         "product_id, date, daily_sales, daily_amount"),
    ]
    for i, (title, filename, count, desc, fields) in enumerate(data_sources):
        x = Inches(0.5 + i * 4.2)
        y = Inches(1.5)
        add_shape(slide, x, y, Inches(3.8), Inches(5.2), RGBColor(0xFF, 0xFF, 0xFF))
        color = [COLORS['secondary'], COLORS['green'], COLORS['orange']][i]
        add_shape(slide, x, y, Inches(3.8), Pt(4), color)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.4), Inches(0.4),
                     title, font_size=20, color=COLORS['primary'], bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.6), Inches(3.4), Inches(0.3),
                     f"文件: {filename}", font_size=12, color=COLORS['muted'])
        add_text_box(slide, x + Inches(0.2), y + Inches(0.9), Inches(3.4), Inches(0.4),
                     count, font_size=28, color=color, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(1.5), Inches(3.4), Inches(1.5),
                     desc, font_size=13, color=COLORS['text'])
        add_text_box(slide, x + Inches(0.2), y + Inches(3.2), Inches(3.4), Inches(0.3),
                     "核心字段:", font_size=12, color=COLORS['primary'], bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(3.5), Inches(3.4), Inches(1.5),
                     fields, font_size=11, color=COLORS['muted'])

    # ================================================================
    # Slide 8: 初始化存储
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_content_slide(slide, "阶段 2：初始化存储 - Backend 容器内")

    add_shape(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(5.5),
              RGBColor(0xFF, 0xFF, 0xFF))
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.4),
                 "HBase 存储", font_size=22, color=COLORS['primary'], bold=True)
    hbase_items = [
        "product 表 (RowKey: product_id)",
        "  → 商品基础信息: name, category, brand, price",
        "review 表 (RowKey: review_id)",
        "  → 评价数据: content, rating, sentiment",
        "monthly_sales 表 (RowKey: product_id_date)",
        "  → 日销量: dailySales, dailyAmount",
    ]
    add_bullet_text(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4),
                    hbase_items, font_size=14, spacing=Pt(6))

    add_shape(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(5.5),
              RGBColor(0xFF, 0xFF, 0xFF))
    add_text_box(slide, Inches(7.1), Inches(1.6), Inches(5), Inches(0.4),
                 "Elasticsearch 索引", font_size=22, color=COLORS['green'], bold=True)
    es_items = [
        "索引名: products",
        "  → name (text, ik_max_word 中文分词)",
        "  → description (text, ik_max_word 中文分词)",
        "  → category, brand (keyword)",
        "  → price, originalPrice (float)",
        "",
        "用途: 全文搜索 + 相似商品推荐",
    ]
    add_bullet_text(slide, Inches(7.1), Inches(2.2), Inches(5.5), Inches(4),
                    es_items, font_size=14, spacing=Pt(6))

    # ================================================================
    # Slide 9: HDFS + Spark 数据清洗
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_content_slide(slide, "阶段 3：HDFS 存储 + Spark 数据清洗")

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(12), Inches(0.4),
                 "HDFS 目录结构", font_size=20, color=COLORS['primary'], bold=True)
    hdfs_paths = [
        "/douyin/raw/products/          → /douyin/cleaned/products/",
        "/douyin/raw/reviews/           → /douyin/cleaned/reviews/",
        "/douyin/raw/monthly_sales/     → /douyin/cleaned/monthly_sales/",
    ]
    add_bullet_text(slide, Inches(0.8), Inches(1.9), Inches(12), Inches(1.5),
                    hdfs_paths, font_size=14, spacing=Pt(8))

    add_text_box(slide, Inches(0.8), Inches(3.5), Inches(12), Inches(0.4),
                 "Spark 清洗规则 (data_cleaning.py)", font_size=20,
                 color=COLORS['primary'], bold=True)

    rules = [
        ("商品数据", "去除首尾空格 | 价格负值设为0 | 过滤空名称/空ID", COLORS['secondary']),
        ("评价数据", "去除首尾空格 | 评分限制1-5 | 过滤空内容/空ID", COLORS['green']),
        ("日销量数据", "销量/金额负值设为0 | 过滤空日期/空ID", COLORS['orange']),
    ]
    for i, (dtype, rule, color) in enumerate(rules):
        y = Inches(4.1 + i * 0.9)
        add_shape(slide, Inches(0.8), y, Pt(4), Inches(0.6), color)
        add_text_box(slide, Inches(1.1), y, Inches(2), Inches(0.4),
                     dtype, font_size=16, color=COLORS['primary'], bold=True)
        add_text_box(slide, Inches(3.5), y, Inches(9), Inches(0.6),
                     rule, font_size=14, color=COLORS['text'])

    # ================================================================
    # Slide 10: Spark 分析计算总览
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_content_slide(slide, "阶段 4：Spark 分析计算（三作业并行）")

    jobs = [
        ("商品分析", "product_analysis_job.py", "EcommerceProductAnalysis",
         "爆款指数(hot_score)\n评价指标(好评率/差评率)\n销量指标(增长率)\njieba关键词提取",
         COLORS['secondary']),
        ("NLP 分析", "nlp_analysis_job.py", "EcommerceNLPAnalysis",
         "jieba中文分词\nSnowNLP情感分析\n高频关键词统计\n情感/评分分布",
         COLORS['green']),
        ("趋势聚合", "trend_aggregation_job.py", "EcommerceTrendAggregation",
         "日销量→月度聚合\n近12月趋势数组\n品类维度聚合\n月度对比分析",
         COLORS['orange']),
    ]
    for i, (title, script, app_name, desc, color) in enumerate(jobs):
        x = Inches(0.5 + i * 4.2)
        y = Inches(1.4)
        add_shape(slide, x, y, Inches(3.8), Inches(5.5), RGBColor(0xFF, 0xFF, 0xFF))
        add_shape(slide, x, y, Inches(3.8), Pt(4), color)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.4), Inches(0.4),
                     title, font_size=22, color=COLORS['primary'], bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.6), Inches(3.4), Inches(0.3),
                     f"脚本: {script}", font_size=11, color=COLORS['muted'])
        add_text_box(slide, x + Inches(0.2), y + Inches(0.9), Inches(3.4), Inches(0.3),
                     f"作业名: {app_name}", font_size=11, color=COLORS['muted'])
        add_shape(slide, x + Inches(0.2), y + Inches(1.3), Inches(3.4), Pt(1), COLORS['bg_light'])
        add_text_box(slide, x + Inches(0.2), y + Inches(1.5), Inches(3.4), Inches(3.5),
                     desc, font_size=15, color=COLORS['text'])

    # ================================================================
    # Slide 11: 商品分析作业详解
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_content_slide(slide, "商品分析作业 - 爆款指数计算")

    add_shape(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.8),
              RGBColor(0xFF, 0xFF, 0xFF))
    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(12), Inches(0.4),
                 "爆款指数公式 (品类内比例法)", font_size=18, color=COLORS['primary'], bold=True)
    formulas = [
        "sales_score  = total_sales / 品类内最大total_sales      (0-1)",
        "growth_score = monthly_growth / 100 (截断到0-1)          (0-1)",
        "rating_score = avg_rating / 5.0                          (0-1)",
        "review_score = review_count / 品类内最大review_count     (0-1)",
        "hot_score = 0.4×sales + 0.3×growth + 0.2×rating + 0.1×review",
    ]
    add_bullet_text(slide, Inches(0.8), Inches(1.9), Inches(12), Inches(1.2),
                    formulas, font_size=13, color=COLORS['text'], spacing=Pt(4))

    add_text_box(slide, Inches(0.8), Inches(3.3), Inches(12), Inches(0.4),
                 "输出到 HBase 的字段", font_size=18, color=COLORS['primary'], bold=True)

    fields = [
        ("metrics:hotScore", "综合爆款指数 (0-1)"),
        ("metrics:monthlyGrowth", "月增长率 (%)"),
        ("metrics:reviewCount", "评价总数"),
        ("metrics:avgRating", "平均评分 (1-5)"),
        ("metrics:positiveRate", "好评率 (评分≥4)"),
        ("metrics:negativeRate", "差评率 (评分≤2)"),
        ("metrics:topTags", "高频关键词 Top 10 (JSON)"),
    ]
    for i, (col, desc) in enumerate(fields):
        row = i // 2
        col_idx = i % 2
        x = Inches(0.8 + col_idx * 6.2)
        y = Inches(3.8 + row * 0.55)
        add_text_box(slide, x, y, Inches(2.8), Inches(0.4),
                     col, font_size=13, color=COLORS['secondary'], bold=True)
        add_text_box(slide, x + Inches(2.8), y, Inches(3), Inches(0.4),
                     desc, font_size=13, color=COLORS['text'])

    # ================================================================
    # Slide 12: NLP 分析 + 趋势聚合
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_content_slide(slide, "NLP 分析 + 趋势聚合")

    add_shape(slide, Inches(0.5), Inches(1.4), Inches(6), Inches(5.5),
              RGBColor(0xFF, 0xFF, 0xFF))
    add_shape(slide, Inches(0.5), Inches(1.4), Inches(6), Pt(4), COLORS['green'])
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.4),
                 "NLP 评价分析", font_size=22, color=COLORS['green'], bold=True)
    nlp_items = [
        "jieba 中文分词 + 停用词过滤",
        "SnowNLP 情感分析 (0-1得分)",
        "  < 0.3 → negative",
        "  0.3-0.7 → neutral",
        "  > 0.7 → positive",
        "",
        "输出到 HBase review_analysis:",
        "  stats:highFreqKeywords (JSON)",
        "  stats:sentimentDistribution",
        "  stats:ratingDistribution",
    ]
    add_bullet_text(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5),
                    nlp_items, font_size=14, spacing=Pt(5))

    add_shape(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(5.5),
              RGBColor(0xFF, 0xFF, 0xFF))
    add_shape(slide, Inches(6.8), Inches(1.4), Inches(6), Pt(4), COLORS['orange'])
    add_text_box(slide, Inches(7.1), Inches(1.6), Inches(5), Inches(0.4),
                 "趋势聚合", font_size=22, color=COLORS['orange'], bold=True)
    trend_items = [
        "日销量 → 月度聚合",
        "(730,000条 → 24,000条月度数据)",
        "",
        "聚合步骤:",
        "  1. 从date提取月份(YYYY-MM)",
        "  2. 按product_id+month聚合",
        "  3. 生成近12月销量数组",
        "  4. 按品类聚合趋势",
        "",
        "输出到 HBase trend_data:",
        "  monthly:totalSales",
        "  monthly:YYYY-MM:sales (动态列)",
    ]
    add_bullet_text(slide, Inches(7.1), Inches(2.2), Inches(5.5), Inches(4.5),
                    trend_items, font_size=14, spacing=Pt(5))

    # ================================================================
    # Slide 13: HDFS 运行截图（新增）
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_content_slide(slide, "HDFS 文件目录 - 运行效果")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(12), Inches(0.4),
                 "Spark 清洗后数据存储在 HDFS /douyin/cleaned/ 和 /douyin/processed/ 目录",
                 font_size=16, color=COLORS['muted'])
    add_screenshot(slide, 'hdfs文件目录.png',
                   Inches(0.3), Inches(1.9), Inches(12.7), Inches(5.2))

    # ================================================================
    # Slide 14: Spark 运行截图（新增）
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_content_slide(slide, "Spark 作业运行日志")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(12), Inches(0.4),
                 "product_analysis_job 在 Spark Standalone 集群上执行，包含 Stage 划分和 Task 并行处理",
                 font_size=16, color=COLORS['muted'])
    add_screenshot(slide, 'spark_product_analysis任务日志.png',
                   Inches(0.3), Inches(1.9), Inches(12.7), Inches(5.2))

    # ================================================================
    # Slide 15: 章节分隔 - 前端展示
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_section_divider(slide, "系统实现与前端展示",
                        "Vue3 + ECharts 构建的数据可视化大屏")

    # ================================================================
    # Slide 16: Dashboard 总览（新增截图）
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_content_slide(slide, "数据概览 Dashboard")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(12), Inches(0.4),
                 "实时掌握电商数据动态：商品总数、评价总数、爆款数量、品类分布、爆款指数趋势",
                 font_size=16, color=COLORS['muted'])
    add_screenshot(slide, '屏幕截图 2026-06-17 225704.png',
                   Inches(0.3), Inches(1.9), Inches(12.7), Inches(5.2))

    # ================================================================
    # Slide 17: 爆款分析（新增截图）
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_content_slide(slide, "数据分析 - 爆款分析")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(12), Inches(0.4),
                 "爆款指数分布、评价评分分布、高增长商品散点图、热门商品排行",
                 font_size=16, color=COLORS['muted'])
    add_screenshot(slide, '921a731facebc2d24b7fdc3af1ff2fcd.png',
                   Inches(0.3), Inches(1.9), Inches(12.7), Inches(5.2))

    # ================================================================
    # Slide 18: 差评分析（新增截图）
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_content_slide(slide, "数据分析 - 差评分析")
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(12), Inches(0.4),
                 "高频关键词 Top10、用户槽点统计、情感分布、评分分布、LDA主题建模、K-Means聚类",
                 font_size=16, color=COLORS['muted'])
    add_screenshot(slide, '屏幕截图 2026-06-17 225757.png',
                   Inches(0.3), Inches(1.9), Inches(12.7), Inches(5.2))

    # ================================================================
    # Slide 19: 卖点推荐 + 相似搜索（新增截图）
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_content_slide(slide, "卖点推荐 + 相似爆款检索")

    # 左侧 - 卖点推荐
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.4),
                 "AI 卖点推荐", font_size=18, color=COLORS['primary'], bold=True)
    add_screenshot(slide, '屏幕截图 2026-06-17 230228.png',
                   Inches(0.3), Inches(1.8), Inches(6.2), Inches(2.8))

    # 右侧 - 相似搜索
    add_text_box(slide, Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.4),
                 "ES 相似爆款检索", font_size=18, color=COLORS['primary'], bold=True)
    add_screenshot(slide, '屏幕截图 2026-06-17 230000.png',
                   Inches(6.8), Inches(1.8), Inches(6.2), Inches(2.8))

    # 底部说明
    add_text_box(slide, Inches(0.8), Inches(5.0), Inches(12), Inches(2),
                 "卖点推荐：基于评价分析生成各品类的质量/客服/外观/物流/功能五大维度卖点建议\n"
                 "相似检索：Elasticsearch ik_max_word 中文分词 + 多维度相似度计算，支持按品类筛选",
                 font_size=14, color=COLORS['text'])

    # ================================================================
    # Slide 20: 关键技术点
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_content_slide(slide, "关键技术点")

    techs = [
        ("品类内比例法", "爆款指数计算",
         "同一品类内，销量最高的商品 sales_score=1.0\n其他商品按比例获得 0-1 值\n确保不同品类之间具有可比性\n\n公式: value / 品类内最大值",
         COLORS['secondary']),
        ("爆款定义", "前端展示逻辑",
         "每个品类按 hot_score 降序排列\n取前 20% 作为爆款\n\nPython: top_count = max(1, int(len(products)*0.2))\n每个品类独立计算，不受其他品类影响",
         COLORS['green']),
        ("HBase 重试", "连接容错机制",
         "装饰器 @_retry_on_disconnect\n连接断开时自动重连并重试\n最多重试 2 次\n\n检测关键词: timed out, connection,\nbroken, socket, reset, closed",
         COLORS['orange']),
    ]
    for i, (title, subtitle, desc, color) in enumerate(techs):
        x = Inches(0.5 + i * 4.2)
        y = Inches(1.4)
        add_shape(slide, x, y, Inches(3.8), Inches(5.5), RGBColor(0xFF, 0xFF, 0xFF))
        add_shape(slide, x, y, Inches(3.8), Pt(4), color)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.4), Inches(0.4),
                     title, font_size=22, color=COLORS['primary'], bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.6), Inches(3.4), Inches(0.3),
                     subtitle, font_size=14, color=COLORS['muted'])
        add_shape(slide, x + Inches(0.2), y + Inches(1), Inches(3.4), Pt(1), COLORS['bg_light'])
        add_text_box(slide, x + Inches(0.2), y + Inches(1.2), Inches(3.4), Inches(4),
                     desc, font_size=13, color=COLORS['text'])

    # ================================================================
    # Slide 21: 数据量级汇总
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_content_slide(slide, "数据量级汇总")

    metrics = [
        ("75.5万", "原始记录", COLORS['secondary']),
        ("6,000", "分析结果", COLORS['green']),
        ("2,000", "ES索引", COLORS['orange']),
        ("6张", "HBase表", COLORS['primary']),
    ]
    for i, (num, label, color) in enumerate(metrics):
        add_metric_card(slide, Inches(0.5 + i * 3.2), Inches(1.4),
                        Inches(2.8), Inches(1.2), num, label, color)

    add_text_box(slide, Inches(0.8), Inches(3), Inches(12), Inches(0.4),
                 "各阶段数据量", font_size=18, color=COLORS['primary'], bold=True)

    table_data = [
        ("Mock 生成", "商品 2,000 | 评价 22,970 | 日销量 730,000"),
        ("HBase 原始", "product 2,000 | review 22,970 | monthly_sales 730,000"),
        ("ES 索引", "products 2,000 (全文搜索)"),
        ("Spark 清洗", "HDFS cleaned 目录 (与原始数据量一致)"),
        ("Spark 分析", "product_analysis 2,000 | review_analysis 2,000 | trend_data 2,000"),
    ]
    for i, (stage, detail) in enumerate(table_data):
        y = Inches(3.5 + i * 0.7)
        add_shape(slide, Inches(0.8), y, Inches(2.5), Inches(0.5), COLORS['primary'])
        add_text_box(slide, Inches(0.8), y + Pt(4), Inches(2.5), Inches(0.4),
                     stage, font_size=13, color=COLORS['light_text'], bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(3.5), y + Pt(4), Inches(9), Inches(0.4),
                     detail, font_size=13, color=COLORS['text'])

    # ================================================================
    # Slide 22: 结论与创新点（新增）
    # ================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide, COLORS['dark'])
    add_shape(slide, 0, 0, SLIDE_WIDTH, Pt(4), COLORS['accent'])

    add_text_box(slide, Inches(1.5), Inches(0.5), Inches(10), Inches(0.8),
                 "结论与创新点", font_size=40, color=COLORS['light_text'],
                 bold=True, alignment=PP_ALIGN.CENTER)

    # 四个创新点
    innovations = [
        ("创新点 1", "品类内比例法爆款指数",
         "采用品类内归一化方法，解决跨品类可比性问题。销量、增长率、评分、评价量四维加权，权重可调。",
         COLORS['secondary']),
        ("创新点 2", "NLP 多维评价分析",
         "jieba分词 + SnowNLP情感分析 + LDA主题建模 + K-Means聚类，四层分析挖掘用户真实反馈。",
         COLORS['green']),
        ("创新点 3", "ES 全文搜索 + 相似推荐",
         "Elasticsearch ik_max_word 中文分词，支持全文检索和多维度相似商品推荐。",
         COLORS['orange']),
        ("创新点 4", "Airflow 全链路自动化",
         "DAG 编排 15 个任务，每日凌晨自动执行完整 ETL 流程，全部 Docker 容器化部署。",
         COLORS['accent']),
    ]

    for i, (label, title, desc, color) in enumerate(innovations):
        x = Inches(0.5 + i * 3.15)
        y = Inches(1.6)
        # 卡片
        add_shape(slide, x, y, Inches(2.9), Inches(3.5), RGBColor(0x25, 0x25, 0x3F))
        add_shape(slide, x, y, Inches(2.9), Pt(4), color)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(2.6), Inches(0.3),
                     label, font_size=12, color=color, bold=True)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.5), Inches(2.6), Inches(0.4),
                     title, font_size=16, color=COLORS['light_text'], bold=True)
        add_text_box(slide, x + Inches(0.15), y + Inches(1.0), Inches(2.6), Inches(2.2),
                     desc, font_size=12, color=COLORS['muted'])

    # 展望
    add_text_box(slide, Inches(0.8), Inches(5.4), Inches(12), Inches(0.4),
                 "未来展望", font_size=20, color=COLORS['accent'], bold=True)
    outlook = [
        "接入真实抖音电商API，实现真实数据采集",
        "引入Flink实时流处理，支持实时销量监控",
        "集成机器学习模型，实现销量预测和智能推荐",
    ]
    add_bullet_text(slide, Inches(0.8), Inches(5.9), Inches(12), Inches(1.5),
                    outlook, font_size=14, color=COLORS['light_text'], spacing=Pt(8))

    # 谢谢
    add_text_box(slide, Inches(1.5), Inches(7.0), Inches(10), Inches(0.4),
                 "谢谢！", font_size=24, color=COLORS['accent'],
                 bold=True, alignment=PP_ALIGN.CENTER)

    # ================================================================
    # 保存
    # ================================================================
    output_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'docs')
    output_path = os.path.join(output_dir, '电商数据洞察平台-项目演示v3.pptx')
    prs.save(output_path)
    print(f"PPT 已生成: {output_path}")
    print(f"共 {len(prs.slides)} 页幻灯片")
    return output_path


if __name__ == "__main__":
    generate_ppt()
